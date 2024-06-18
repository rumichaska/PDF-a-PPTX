""" Importación de módulos
"""

import os
import fnmatch
import fitz
import cv2
import numpy as np

from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas
from PIL import Image

# FUNCIONES ----

def pdf_to_images(pdf_path, filename, zoom_factor=2):
    """Covierte un archivo pdf en imágenes"""
    # Patrón del nombre del pdf
    filename_pattern = filename[0:len(filename) - 4]
    # Crear directorio de imágenes de las diapositivas
    page_dir = os.path.exists("./pages/")
    if not page_dir:
        os.makedirs("./pages/")
    # Crear subdirectorio del pdf a convertir
    pdf_slide_dir = os.path.exists(f"./pages/{filename_pattern}")
    if not pdf_slide_dir:
        os.makedirs(f"./pages/{filename_pattern}")

    # Abre el archivo PDF
    pdf_document = fitz.open(f"{pdf_path}{filename}")
    # Iteración por páginas
    for page_number in range(len(pdf_document)):
        # Extrae la página como imagen con mayor resolución
        page = pdf_document.load_page(page_number)
        matrix = fitz.Matrix(zoom_factor, zoom_factor)
        pix = page.get_pixmap(matrix=matrix)
        # Guardar temporalmente la imágen de la página
        new_page_number = page_number + 1
        slide_number = f"0{new_page_number}" if new_page_number < 10 else f"{new_page_number}"
        img_path = f"./pages/{filename_pattern}/page_{slide_number}.png"
        pix.save(img_path)


def has_purple_border(image_path):
    """Verifica si la imágen tiene borde púrpura y retorna las coordenadas"""
    # Cargar imágen
    new_image = cv2.imread(image_path)
    # Convertir imágen a formato HSV para detectar color patrón
    hsv = cv2.cvtColor(new_image, cv2.COLOR_BGR2HSV)
    # Definir el rango de colores HSV patrón
    lower_purple = np.array([120, 40, 40])
    upper_purple = np.array([150, 255, 255])
    # Crear unas máscara para el color patrón
    mask_purple = cv2.inRange(hsv, lower_purple, upper_purple)
    # Ubicar contornos de la máscara
    contours_purple, _ = cv2.findContours(mask_purple, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    # Descartar pequeños contornos que no sean del color patrón
    contours_purple = [cnt for cnt in contours_purple if cv2.contourArea(cnt) > 100]
    return contours_purple


def get_content(png_path, out_path, filename):
    """Extrae información a partir de imágenes con borde púrpura"""
    # Crear directorio de imágenes de las diapositivas
    table_dir = os.path.exists(out_path)
    if not table_dir:
        os.makedirs(out_path)

    # Validar si la imágen tiene borde púrpura
    new_image = cv2.imread(f"{png_path}{filename}")
    image_contours = has_purple_border(f"{png_path}{filename}")
    contours_validation = len(image_contours) > 0
    # Extraer información cuando detecta color patrón
    if contours_validation:
        # Asumir que el borde patrón mas grande tiene el contenido de interés
        largest_contour_purple = max(image_contours, key=cv2.contourArea)
        x, y, w, h = cv2.boundingRect(largest_contour_purple)
        # Ajustar las coordenadas para excluir el borde patrón
        # NOTE: Ajustar este margen según las necesidades
        border_margin = 12
        x = max(0, x + border_margin)
        y = max(0, y + border_margin)
        w = max(0, w - 2 * border_margin)
        h = max(0, h - 2 * border_margin)
        # Cortar la imágen en función al recuadro del borde patrón
        table_image_new = new_image[y:y+h, x:x+w]
        # Guardar el contenido de interés
        cropped_table_path = f"{out_path}{filename}"
        cv2.imwrite(cropped_table_path, table_image_new)


def images_to_pptx(images, image_path, pptx_file):
    """Convierte imágenes png a un archivo pptx"""
    # Crea una presentación en blanco
    presentation = Presentation()
    # Define el tamaño de la imagen y la posición
    left = top = Inches(0)
    # Generar diapositivas con imágenes del pdf
    for image in images:
        # Añade una diapositiva en blanco a la presentación
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        # Añadir la imágena a la diapositiva creada
        slide.shapes.add_picture(f"{image_path}{image}", left, top, height=Inches(7.5))
    # Guarda la presentación
    presentation.save(f"{pptx_file}.pptx")


def images_to_pdf(images, image_path, pdf_file):
    """Convierte imágenes png a un archivo pdf"""
    # Primera imágen
    first_image = Image.open(f"{image_path}{images[0]}")
    # Preparar canvas para generación de pdf
    c = canvas.Canvas(f"{pdf_file}.pdf", pagesize=(first_image.width, first_image.height))
    for image in images:
        image_file = f"{image_path}{image}"
        image = Image.open(image_file)
        c.setPageSize((image.width, image.height))
        c.drawImage(image_file, 0, 0, image.width, image.height)
        c.showPage()
    c.save()

# CONVERTIR SALAS DE PDF A PPTX ----

# Listar archivos de las salas
list_files = fnmatch.filter(os.listdir("./content/"), "*.pdf")

# Generación de salas
# # NOTE: Ajustar el `zoom_factor` para cambiar la resolución
for file in list_files:
    # Proyecto
    project_name = file[0:len(file) - 4]
    # Convertir archivos
    pdf_to_images("./content/", file, zoom_factor=3)
    # Directorio del archivo convertido
    dir_in = f"./pages/{project_name}/"
    dir_out = f"./content/{project_name}/"

    # Diapositiva
    list_images = fnmatch.filter(os.listdir(dir_in), "*.png")
    # Extraer tablas
    for table in list_images:
        get_content(dir_in, dir_out, table)

    # Imágenes a exportar
    list_images_tables = fnmatch.filter(os.listdir(dir_out), "*.png")
    list_out = list(set(list_images) - set(list_images_tables))
    list_out.sort()
    # Generar pptx
    images_to_pptx(list_out, dir_in, f"{dir_out}{project_name}")
    # Generar pdf
    images_to_pdf(list_out, dir_in, f"{dir_out}{project_name}")
